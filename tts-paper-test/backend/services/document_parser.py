"""
多格式文档解析服务
支持: docx, doc, xlsx, xls, pdf, pptx, ppt, md, txt, csv, xmind
"""
import os
import csv
import io
from pathlib import Path
from typing import Optional

from core.logging_config import get_logger

logger = get_logger("document_parser")

# 支持的扩展名
SUPPORTED_EXTENSIONS = {
    ".docx", ".doc", ".xlsx", ".xls", ".pdf",
    ".pptx", ".ppt", ".md", ".txt", ".csv", ".xmind",
}

SUPPORTED_EXTENSIONS_LABEL = ", ".join(sorted(SUPPORTED_EXTENSIONS))


class DocumentParseError(Exception):
    """文档解析异常"""
    pass


def parse_document(file_path: str) -> str:
    """
    解析文档并提取文本内容
    返回提取的纯文本
    """
    ext = Path(file_path).suffix.lower()
    logger.info(f"开始解析文档: {file_path}, 类型: {ext}")

    if ext == ".docx":
        return _parse_docx(file_path)
    elif ext == ".doc":
        return _parse_doc(file_path)
    elif ext == ".xlsx":
        return _parse_xlsx(file_path)
    elif ext == ".xls":
        return _parse_xls(file_path)
    elif ext == ".pdf":
        return _parse_pdf(file_path)
    elif ext == ".pptx":
        return _parse_pptx(file_path)
    elif ext == ".ppt":
        return _parse_ppt(file_path)
    elif ext == ".md":
        return _parse_text(file_path)
    elif ext == ".txt":
        return _parse_text(file_path)
    elif ext == ".csv":
        return _parse_csv(file_path)
    elif ext == ".xmind":
        return _parse_xmind(file_path)
    else:
        raise DocumentParseError(f"不支持的文件格式: {ext}")


def _parse_docx(file_path: str) -> str:
    """解析 .docx 文件"""
    try:
        from docx import Document
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # 也提取表格内容
        tables_text = []
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    tables_text.append(" | ".join(cells))
        content = "\n".join(paragraphs)
        if tables_text:
            content += "\n\n--- 表格内容 ---\n" + "\n".join(tables_text)
        return content
    except ImportError:
        raise DocumentParseError("python-docx 未安装，无法解析 .docx 文件")
    except Exception as e:
        raise DocumentParseError(f"解析 .docx 失败: {e}")


def _parse_doc(file_path: str) -> str:
    """尝试解析旧版 .doc 文件"""
    # 优先使用 antiword / catdoc 系统命令
    import subprocess
    try:
        result = subprocess.run(["antiword", file_path], capture_output=True, text=True, timeout=30)
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
    except FileNotFoundError:
        pass
    try:
        result = subprocess.run(["catdoc", file_path], capture_output=True, text=True, timeout=30)
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
    except FileNotFoundError:
        pass
    # 尝试用 python-docx 打开 (有时能兼容)
    try:
        return _parse_docx(file_path)
    except Exception:
        raise DocumentParseError(
            "无法解析 .doc 文件，请安装 antiword 或 catdoc:\n"
            "  macOS: brew install antiword\n"
            "  Linux: apt install antiword\n"
            "或转换为 .docx 格式后重试"
        )


def _parse_xlsx(file_path: str) -> str:
    """解析 .xlsx 文件"""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        lines = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            lines.append(f"\n## Sheet: {sheet_name}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    lines.append(" | ".join(cells))
        wb.close()
        return "\n".join(lines)
    except ImportError:
        raise DocumentParseError("openpyxl 未安装，无法解析 .xlsx 文件")
    except Exception as e:
        raise DocumentParseError(f"解析 .xlsx 失败: {e}")


def _parse_xls(file_path: str) -> str:
    """解析 .xls 文件"""
    try:
        import xlrd
        wb = xlrd.open_workbook(file_path)
        lines = []
        for sheet_name in wb.sheet_names():
            ws = wb.sheet_by_name(sheet_name)
            lines.append(f"\n## Sheet: {sheet_name}")
            for row_idx in range(ws.nrows):
                cells = [str(ws.cell_value(row_idx, c)).strip() for c in range(ws.ncols)]
                cells = [c for c in cells if c]
                if cells:
                    lines.append(" | ".join(cells))
        return "\n".join(lines)
    except ImportError:
        raise DocumentParseError("xlrd 未安装，无法解析 .xls 文件")
    except Exception as e:
        raise DocumentParseError(f"解析 .xls 失败: {e}")


def _parse_pdf(file_path: str) -> str:
    """解析 .pdf 文件"""
    try:
        import pdfplumber
        texts = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text and text.strip():
                    texts.append(text.strip())
                # 也提取表格
                tables = page.extract_tables()
                for table in tables:
                    rows_text = []
                    for row in table:
                        cells = [str(c).strip() for c in row if c and str(c).strip()]
                        if cells:
                            rows_text.append(" | ".join(cells))
                    if rows_text:
                        texts.append("--- 表格 ---\n" + "\n".join(rows_text))
        return "\n\n".join(texts) if texts else ""
    except ImportError:
        raise DocumentParseError("pdfplumber 未安装，无法解析 .pdf 文件")
    except Exception as e:
        raise DocumentParseError(f"解析 .pdf 失败: {e}")


def _parse_pptx(file_path: str) -> str:
    """解析 .pptx 文件"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_lines = [f"## Slide {i}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_lines.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            slide_lines.append(" | ".join(cells))
            slides_text.append("\n".join(slide_lines))
        return "\n\n".join(slides_text)
    except ImportError:
        raise DocumentParseError("python-pptx 未安装，无法解析 .pptx 文件")
    except Exception as e:
        raise DocumentParseError(f"解析 .pptx 失败: {e}")


def _parse_ppt(file_path: str) -> str:
    """尝试解析旧版 .ppt 文件"""
    import subprocess
    try:
        result = subprocess.run(["catppt", file_path], capture_output=True, text=True, timeout=30)
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
    except FileNotFoundError:
        pass
    raise DocumentParseError(
        "无法解析 .ppt 文件，请转换为 .pptx 格式后重试"
    )


def _parse_text(file_path: str) -> str:
    """解析纯文本 / Markdown 文件"""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        # 尝试其他编码
        for enc in ["gbk", "gb2312", "latin-1"]:
            try:
                with open(file_path, "r", encoding=enc) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        raise DocumentParseError("无法识别文本编码")


def _parse_csv(file_path: str) -> str:
    """解析 .csv 文件"""
    try:
        lines = []
        with open(file_path, "r", encoding="utf-8") as f:
            reader = csv.reader(f)
            for row in reader:
                cells = [c.strip() for c in row if c.strip()]
                if cells:
                    lines.append(" | ".join(cells))
        return "\n".join(lines) if lines else ""
    except UnicodeDecodeError:
        try:
            lines = []
            with open(file_path, "r", encoding="gbk") as f:
                reader = csv.reader(f)
                for row in reader:
                    cells = [c.strip() for c in row if c.strip()]
                    if cells:
                        lines.append(" | ".join(cells))
            return "\n".join(lines) if lines else ""
        except Exception as e:
            raise DocumentParseError(f"解析 .csv 失败: {e}")
    except Exception as e:
        raise DocumentParseError(f"解析 .csv 失败: {e}")


def _parse_xmind(file_path: str) -> str:
    """解析 .xmind 文件"""
    try:
        import xmindparser
        result = xmindparser.xmind_to_dict(file_path)
        lines = []
        def extract_topic(topic, depth=0):
            """递归提取主题"""
            prefix = "  " * depth
            title = topic.get("title", "")
            if title:
                lines.append(f"{prefix}- {title}")
            # 提取备注
            note = topic.get("note", "")
            if note:
                lines.append(f"{prefix}  > {note}")
            # 子主题
            for child in topic.get("topics", []):
                extract_topic(child, depth + 1)

        for sheet in result:
            sheet_title = sheet.get("title", "思维导图")
            lines.append(f"# {sheet_title}")
            topic = sheet.get("topic", {})
            extract_topic(topic)
        return "\n".join(lines)
    except ImportError:
        raise DocumentParseError("xmindparser 未安装，无法解析 .xmind 文件")
    except Exception as e:
        raise DocumentParseError(f"解析 .xmind 失败: {e}")
